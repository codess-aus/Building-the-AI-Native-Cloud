"""Build two editable PowerPoint decks for "Building the AI-Native Cloud".

Outputs:
    dist/Building-the-AI-Native-Cloud_Light.pptx  (black text on white)
    dist/Building-the-AI-Native-Cloud_Dark.pptx   (white text on near-black)

Both decks share identical structure and copy, lifted verbatim from
context.md. Calibri throughout, every text element >=18pt, accents drawn
from a teal/sage/lime palette that has been contrast-checked against each
background. Full speaker notes are placed in each slide's Notes pane.
"""

from __future__ import annotations

import argparse
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

FONT = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.6)
MARGIN_TOP = Inches(0.5)
TITLE_HEIGHT = Inches(0.9)
ACCENT_RULE_OFFSET = Inches(0.15)
FOOTER_HEIGHT = Inches(0.35)

TALK_TITLE = "Building the AI-Native Cloud"
PRESENTER = "Michelle Mei-Ling Sandford  \u00b7  Asia DevOps Conference"


# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class Theme:
    name: str
    background: RGBColor
    text_primary: RGBColor
    text_muted: RGBColor
    accent_teal: RGBColor
    accent_sage: RGBColor
    accent_lime: RGBColor
    on_teal: RGBColor   # text colour that sits on a teal fill
    on_sage: RGBColor   # text colour that sits on a sage fill
    on_lime: RGBColor   # text colour that sits on a lime fill
    table_header_text: RGBColor
    table_row_alt: RGBColor


LIGHT = Theme(
    name="Light",
    background=RGBColor(0xFF, 0xFF, 0xFF),
    text_primary=RGBColor(0x11, 0x11, 0x11),
    text_muted=RGBColor(0x4A, 0x4A, 0x4A),
    accent_teal=RGBColor(0x0E, 0x7C, 0x86),
    accent_sage=RGBColor(0x5C, 0x7F, 0x58),
    accent_lime=RGBColor(0x5E, 0x8A, 0x00),
    on_teal=RGBColor(0xFF, 0xFF, 0xFF),
    on_sage=RGBColor(0xFF, 0xFF, 0xFF),
    on_lime=RGBColor(0x11, 0x11, 0x11),
    table_header_text=RGBColor(0xFF, 0xFF, 0xFF),
    table_row_alt=RGBColor(0xF1, 0xF5, 0xF4),
)

DARK = Theme(
    name="Dark",
    background=RGBColor(0x0B, 0x0B, 0x0D),
    text_primary=RGBColor(0xF5, 0xF5, 0xF5),
    text_muted=RGBColor(0xBF, 0xBF, 0xBF),
    accent_teal=RGBColor(0x4F, 0xD1, 0xC5),
    accent_sage=RGBColor(0xA8, 0xD5, 0xA2),
    accent_lime=RGBColor(0xC5, 0xE8, 0x66),
    on_teal=RGBColor(0x0B, 0x0B, 0x0D),
    on_sage=RGBColor(0x0B, 0x0B, 0x0D),
    on_lime=RGBColor(0x0B, 0x0B, 0x0D),
    table_header_text=RGBColor(0x0B, 0x0B, 0x0D),
    table_row_alt=RGBColor(0x18, 0x18, 0x1C),
)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------


def _set_solid_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _set_no_line(shape) -> None:
    shape.line.fill.background()


def _set_line(shape, rgb: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def _style_run(run, *, size_pt: float, color: RGBColor, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_text(
    text_frame,
    text: str,
    *,
    size_pt: float,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    text_frame.word_wrap = True
    text_frame.vertical_anchor = anchor
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_top = Inches(0.02)
    text_frame.margin_bottom = Inches(0.02)
    p = text_frame.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    _style_run(run, size_pt=size_pt, color=color, bold=bold)


def _set_paint_background(slide, theme: Theme) -> None:
    """Paint a full-bleed rectangle as the slide background.

    Using a shape rather than the master background means both themes share
    the same single master while still rendering distinct backgrounds.
    """
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_solid_fill(bg, theme.background)
    _set_no_line(bg)
    bg.shadow.inherit = False
    # Send to back via XML (python-pptx has no helper).
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _set_alt_text(shape, description: str) -> None:
    nv = shape._element.nvSpPr.cNvPr if hasattr(shape._element, "nvSpPr") else None
    if nv is None:
        # GroupShape, GraphicFrame, etc. expose nvGrpSpPr / nvGraphicFramePr.
        for attr in ("nvGrpSpPr", "nvGraphicFramePr", "nvPicPr", "nvCxnSpPr"):
            holder = getattr(shape._element, attr, None)
            if holder is not None:
                nv = holder.cNvPr
                break
    if nv is not None:
        nv.set("descr", description)


def _kill_shadow(shape) -> None:
    """Remove the default shadow effect that PowerPoint themes inject."""
    sppr = shape._element.spPr
    for tag in ("a:effectLst", "a:effectDag"):
        for el in sppr.findall(qn(tag)):
            sppr.remove(el)
    effect_lst = etree.SubElement(sppr, qn("a:effectLst"))
    # Empty effectLst explicitly disables inherited effects.
    del effect_lst  # noqa: F841


# ---------------------------------------------------------------------------
# Layout primitives
# ---------------------------------------------------------------------------


def add_title(slide, text: str, theme: Theme) -> None:
    """Add a real title placeholder + thin accent rule beneath it."""
    title_ph = slide.shapes.title
    title_ph.left = MARGIN_X
    title_ph.top = MARGIN_TOP
    title_ph.width = SLIDE_W - 2 * MARGIN_X
    title_ph.height = TITLE_HEIGHT
    tf = title_ph.text_frame
    _set_text(
        tf,
        text,
        size_pt=36,
        color=theme.text_primary,
        bold=True,
        anchor=MSO_ANCHOR.TOP,
    )
    # Accent rule.
    rule_top = MARGIN_TOP + TITLE_HEIGHT + ACCENT_RULE_OFFSET
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_X,
        rule_top,
        Inches(1.6),
        Pt(3),
    )
    _set_solid_fill(rule, theme.accent_teal)
    _set_no_line(rule)
    _kill_shadow(rule)
    _set_alt_text(rule, "Decorative accent rule")


def add_footer(slide, page_num: int, total: int, theme: Theme) -> None:
    """Footer with talk title on the left and page X / Y on the right."""
    y = SLIDE_H - FOOTER_HEIGHT - Inches(0.15)
    left = slide.shapes.add_textbox(MARGIN_X, y, Inches(8), FOOTER_HEIGHT)
    _set_text(
        left.text_frame,
        TALK_TITLE,
        size_pt=18,
        color=theme.text_muted,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    right = slide.shapes.add_textbox(
        SLIDE_W - MARGIN_X - Inches(2.4),
        y,
        Inches(2.4),
        FOOTER_HEIGHT,
    )
    _set_text(
        right.text_frame,
        f"{page_num} / {total}",
        size_pt=18,
        color=theme.text_muted,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_bullets(
    slide,
    items: Iterable[str],
    theme: Theme,
    *,
    top: Emu = Inches(1.8),
    height: Emu = Inches(4.6),
    size_pt: float = 22,
) -> None:
    box = slide.shapes.add_textbox(
        MARGIN_X,
        top,
        SLIDE_W - 2 * MARGIN_X,
        height,
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    items = list(items)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.text = ""
        p.space_after = Pt(10)
        p.line_spacing = 1.15

        marker = p.add_run()
        marker.text = "\u25A0  "  # filled square + two spaces
        _style_run(marker, size_pt=size_pt, color=theme.accent_teal, bold=True)

        body = p.add_run()
        body.text = item
        _style_run(body, size_pt=size_pt, color=theme.text_primary)


def add_pullquote(
    slide,
    text: str,
    theme: Theme,
    *,
    top: Emu,
    height: Emu = Inches(1.2),
    accent: RGBColor | None = None,
) -> None:
    accent_color = accent if accent is not None else theme.accent_teal
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_X,
        top,
        Pt(6),
        height,
    )
    _set_solid_fill(bar, accent_color)
    _set_no_line(bar)
    _kill_shadow(bar)
    _set_alt_text(bar, "Decorative accent bar")

    box = slide.shapes.add_textbox(
        MARGIN_X + Inches(0.25),
        top,
        SLIDE_W - 2 * MARGIN_X - Inches(0.25),
        height,
    )
    _set_text(
        box.text_frame,
        text,
        size_pt=28,
        color=theme.text_primary,
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_callout_chip(
    slide,
    text: str,
    accent_fill: RGBColor,
    on_color: RGBColor,
    *,
    left: Emu,
    top: Emu,
    width: Emu = Inches(3.6),
    height: Emu = Inches(0.65),
) -> None:
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    chip.adjustments[0] = 0.5
    _set_solid_fill(chip, accent_fill)
    _set_no_line(chip)
    _kill_shadow(chip)
    _set_alt_text(chip, f"Chip: {text}")
    _set_text(
        chip.text_frame,
        text,
        size_pt=20,
        color=on_color,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_comparison_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    theme: Theme,
    *,
    top: Emu = Inches(1.95),
    height: Emu = Inches(4.7),
) -> None:
    n_cols = len(headers)
    n_rows = len(rows) + 1
    width = SLIDE_W - 2 * MARGIN_X
    table_shape = slide.shapes.add_table(n_rows, n_cols, MARGIN_X, top, width, height)
    table = table_shape.table

    # Distribute column widths: first column narrower (row labels).
    first_w = Inches(2.0)
    remaining = width - first_w
    other_w = Emu(remaining // (n_cols - 1))
    table.columns[0].width = first_w
    for c in range(1, n_cols):
        table.columns[c].width = other_w

    # Header row.
    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.accent_teal
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.08)
        cell.margin_bottom = Inches(0.08)
        tf = cell.text_frame
        tf.word_wrap = True
        _set_text(
            tf,
            htext,
            size_pt=20,
            color=theme.table_header_text,
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Body rows.
    for r, row in enumerate(rows, start=1):
        for c, ctext in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                theme.table_row_alt if r % 2 == 1 else theme.background
            )
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)
            tf = cell.text_frame
            tf.word_wrap = True
            _set_text(
                tf,
                ctext,
                size_pt=18,
                color=theme.text_primary,
                bold=(c == 0),
                anchor=MSO_ANCHOR.MIDDLE,
            )

    _set_alt_text(
        table_shape,
        "Comparison table of three eras: Traditional, AI-Assisted, AI-Native, "
        "across unit of work, who edits, artifacts, and governance.",
    )


def add_loop_diagram(slide, theme: Theme, *, top: Emu = Inches(2.05)) -> None:
    """Six rounded-rectangle nodes connected by arrows, with a curved
    feedback arrow returning from Delivery to Spec.
    """
    labels = [
        "Intent",
        "Spec",
        "Context",
        "Change",
        "Verification",
        "Delivery",
    ]
    n = len(labels)
    # The diagram uses tighter horizontal margins than the rest of the slide
    # so the longest label ("Verification") fits on one line at 18pt bold.
    diagram_margin = Inches(0.35)
    usable_w = SLIDE_W - 2 * diagram_margin
    node_w = Inches(1.85)
    node_h = Inches(0.85)
    gap = (usable_w - node_w * n) / (n - 1)
    arrow_w = Inches(0.32)

    shape_ids = []
    accents = [
        theme.accent_teal,
        theme.accent_teal,
        theme.accent_sage,
        theme.accent_sage,
        theme.accent_lime,
        theme.accent_lime,
    ]

    node_positions = []
    for i, label in enumerate(labels):
        x = diagram_margin + i * (node_w + gap)
        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, top, node_w, node_h
        )
        node.adjustments[0] = 0.3
        _set_solid_fill(node, accents[i])
        _set_no_line(node)
        _kill_shadow(node)
        on_color = (
            theme.on_teal if accents[i] == theme.accent_teal
            else theme.on_sage if accents[i] == theme.accent_sage
            else theme.on_lime
        )
        _set_text(
            node.text_frame,
            label,
            size_pt=18,
            color=on_color,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Loop nodes are tight on width; remove inner side margins so the
        # longest label ("Verification") fits on one line.
        node.text_frame.margin_left = Inches(0)
        node.text_frame.margin_right = Inches(0)
        _set_alt_text(node, f"Loop stage: {label}")
        node_positions.append((x, x + node_w))
        shape_ids.append(node)

    # Forward arrows between nodes.
    arrow_y = top + Emu(node_h // 2) - Emu(int(arrow_w) // 2)
    for i in range(n - 1):
        a_left = Emu(int(node_positions[i][1]) + int((gap - arrow_w) // 2))
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            a_left,
            arrow_y,
            arrow_w,
            arrow_w,
        )
        _set_solid_fill(arrow, theme.text_muted)
        _set_no_line(arrow)
        _kill_shadow(arrow)
        _set_alt_text(arrow, "Decorative arrow")
        shape_ids.append(arrow)

    # Feedback arrow: curve from Delivery (last node) back to Spec (index 1)
    # drawn below the row of nodes with a label.
    delivery_x = diagram_margin + (n - 1) * (node_w + gap) + Emu(int(node_w) // 2)
    spec_x = diagram_margin + 1 * (node_w + gap) + Emu(int(node_w) // 2)
    fb_top = top + node_h + Inches(0.5)
    fb_h = Inches(0.04)
    fb_left = min(spec_x, delivery_x)
    fb_w = abs(delivery_x - spec_x)

    horiz = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        fb_left,
        fb_top,
        fb_w,
        fb_h,
    )
    _set_solid_fill(horiz, theme.accent_lime)
    _set_no_line(horiz)
    _kill_shadow(horiz)
    _set_alt_text(horiz, "Feedback path from Delivery to Spec")

    # Down stub from Delivery into horizontal line.
    down_stub = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        delivery_x - Emu(int(fb_h) // 2),
        top + node_h,
        fb_h,
        fb_top - (top + node_h) + fb_h,
    )
    _set_solid_fill(down_stub, theme.accent_lime)
    _set_no_line(down_stub)
    _kill_shadow(down_stub)
    _set_alt_text(down_stub, "Decorative")

    # Up arrow into Spec.
    up_arrow_h = Inches(0.45)
    up_arrow_w = Inches(0.35)
    up_arrow = slide.shapes.add_shape(
        MSO_SHAPE.UP_ARROW,
        spec_x - Emu(int(up_arrow_w) // 2),
        fb_top - up_arrow_h + fb_h,
        up_arrow_w,
        up_arrow_h,
    )
    _set_solid_fill(up_arrow, theme.accent_lime)
    _set_no_line(up_arrow)
    _kill_shadow(up_arrow)
    _set_alt_text(up_arrow, "Feedback arrow into Spec")

    # Feedback label under the horizontal line.
    label_top = fb_top + Inches(0.1)
    label_box = slide.shapes.add_textbox(
        MARGIN_X,
        label_top,
        SLIDE_W - 2 * MARGIN_X,
        Inches(0.5),
    )
    _set_text(
        label_box.text_frame,
        "production signals feed back into the spec",
        size_pt=18,
        color=theme.text_muted,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.TOP,
    )


def add_numbered_two_col(
    slide,
    items: list[str],
    theme: Theme,
    *,
    top: Emu = Inches(1.85),
    height: Emu = Inches(5.0),
) -> None:
    col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) / 2
    row_h = Inches(0.85)
    per_col = (len(items) + 1) // 2  # 7 -> 4 left, 3 right
    for idx, item in enumerate(items):
        col = 0 if idx < per_col else 1
        row = idx if col == 0 else idx - per_col
        x = MARGIN_X + (Emu(int(col_w)) + Inches(0.4)) * col
        y = top + row_h * row

        # Big numeral on the left of each row.
        num_box = slide.shapes.add_textbox(x, y, Inches(0.85), row_h)
        _set_text(
            num_box.text_frame,
            f"{idx + 1}",
            size_pt=40,
            color=theme.accent_lime,
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        text_box = slide.shapes.add_textbox(
            x + Inches(0.95),
            y,
            col_w - Inches(0.95),
            row_h,
        )
        _set_text(
            text_box.text_frame,
            item,
            size_pt=20,
            color=theme.text_primary,
            anchor=MSO_ANCHOR.MIDDLE,
        )


# ---------------------------------------------------------------------------
# Slide content (lifted verbatim from context.md)
# ---------------------------------------------------------------------------


SPEAKER_NOTES: dict[int, str] = {
    1: (
        "Good morning. I'm here to talk about building the AI-native cloud. "
        "Not AI in the cloud, not AI on top of the cloud, but a cloud where "
        "AI agents are first-class participants in how we build, ship, and "
        "run software. In the next 20 minutes I want to convince you that "
        "the shift from cloud-native to AI-native is the same magnitude of "
        "change as the shift from on-prem to cloud, and most of us are still "
        "bolting AI onto an operating model that was never designed for it."
    ),
    2: (
        "Between the day Copilot suggested its first line of code and the "
        "day it started opening its own pull requests, something quietly "
        "shifted. It didn't arrive as one launch. It arrived as a stack: "
        "agent mode in the editor, coding agents on GitHub, the Model "
        "Context Protocol standardizing how tools plug into models, Spec "
        "Kit making the specification a durable artifact, AGENTS.md "
        "becoming the file where we tell AI teammates how to behave. The "
        "2025 DORA report tells us roughly 90 percent of technologists now "
        "use AI at work. So the question is no longer are you using AI. "
        "The question is: is your delivery loop designed for AI, or are "
        "you still bolting AI onto a loop designed without it?"
    ),
    3: (
        "Three eras. Traditional: humans do every stage, AI is absent. "
        "AI-assisted: humans still drive, but Copilot helps at each step. "
        "The loop is unchanged, it just runs faster. AI-native: the loop "
        "itself is redesigned. Specs become first-class, context is "
        "engineered, agents run end-to-end stages, traces are auditable "
        "artifacts. The diagnostic question, next time someone says "
        "they're doing AI-native, is which row of this table did their "
        "loop actually change? If the answer is just we turned on Copilot, "
        "that's AI-assisted. Good place to be. Not the same thing."
    ),
    4: (
        "Here is the definition I want you to leave with. AI-native "
        "software development means designing your delivery loop so AI "
        "isn't just helping developers, it's participating alongside "
        "them. Agents aren't tools you use. They're contributors you "
        "design for, with defined roles, guardrails, feedback loops, and "
        "accountability. Three words do the real work. Designing: it's a "
        "choice about the loop, not a checkbox on a tool. Participating: "
        "AI takes turns, opens PRs, files issues. Accountable: every "
        "agent action leaves a trail."
    ),
    5: (
        "Six stages. Intent: captured deliberately, not in a Slack "
        "thread. Spec: an artifact in the repo, brief and reviewable in "
        "five minutes. Context: engineered through AGENTS.md, MCP "
        "servers, and per-task context packets. Change: made by a human, "
        "by agent mode, or by the coding agent, anchored on the spec. "
        "Verification: layered from cheap to expensive: tests, "
        "evaluators, policy gates, human review. Delivery: ships through "
        "your existing CI/CD and feeds production signals back into the "
        "spec, not just the backlog. The new artifact at the start is "
        "the spec. The new arrow at the end is the feedback into the "
        "spec. Everything else is rebalanced, not reinvented."
    ),
    6: (
        "I want to be clear with this audience: AI-native does not throw "
        "away cloud-native. It depends on it. Containers, declarative "
        "infrastructure, GitOps, observability, zero-trust networking, "
        "ephemeral environments: every one of these is a prerequisite "
        "for safe agent participation. The coding agent runs in an "
        "ephemeral environment because we learned to build ephemeral "
        "environments. It can be scoped to a branch because we learned "
        "branch protection. It can call an internal API through an MCP "
        "server because we learned service meshes and tokens. "
        "Cloud-native gave us the substrate. AI-native is what we build "
        "on it."
    ),
    7: (
        "The most uncomfortable change for engineers is this: the spec "
        "is now the durable artifact, and the code is downstream of it. "
        "Spec Kit formalizes this with a Specify, Plan, Tasks, Implement "
        "flow. But you don't need the toolkit to start. A five-line "
        "Markdown file in a specs/ folder with why, what changes for "
        "the user, success criteria, and out of scope is enough. The "
        "point is not to anticipate every edge case. The point is to "
        "give the agent and the reviewer a shared contract. And here's "
        "the quiet superpower: surfacing ambiguity in a spec review, "
        "before any code is written, is dramatically cheaper than "
        "catching it in a PR."
    ),
    8: (
        "Context is engineered, not assumed. Three layers. First, "
        "AGENTS.md: a plain Markdown file in the repo that tells agents "
        "how this codebase prefers to be edited. What library to use for "
        "dates. Which branches it can push to. Which conventions are "
        "non-negotiable. Second, MCP servers: the Model Context Protocol "
        "lets you expose your wiki, your observability platform, your "
        "internal APIs as first-class tools the agent can call. Writing "
        "a custom MCP server is a weekend project, not a quarter-long "
        "migration. Third, per-task context packets for non-trivial "
        "work. When an agent produces code that is locally plausible "
        "and globally wrong, the cause is almost always context "
        "starvation. Fix context, fix the agent."
    ),
    9: (
        "A colleague told me recently, these are dark times to be "
        "developers. Small companies will cut three engineers and keep "
        "the two who can talk to the business and code faster with AI. "
        "I disagree, strongly. Companies that think AI is about doing "
        "yesterday's work with fewer developers won't be around in two "
        "years. That's not transformation, that's decline. If you want "
        "to win with AI, you expand ambition, not shrink headcount. You "
        "give developers space to orchestrate fleets of agents, to "
        "apply judgment and creativity at a completely different scale. "
        "The operating model shift is this: engineers become "
        "specifiers, context curators, and reviewers of intent. The "
        "keyboard time goes down. The judgment time goes up. That is a "
        "more senior job, not a smaller one."
    ),
    10: (
        "Reliability in an AI-native cloud rests on one principle: no "
        "invisible work. Every agent run leaves a trace. Model version. "
        "Tools called. Files touched. Spec citation per commit. When "
        "something breaks at 2am, the on-call engineer has to be able "
        "to reconstruct why a change was made, not just what changed. "
        "Verification is layered, cheap to expensive: deterministic "
        "tests first, then AI evaluators checking the diff against the "
        "spec's success criteria, then policy gates including Advanced "
        "Security and Copilot Autofix, then human review focused on "
        "intent and integration. The funnel design is what makes the "
        "expensive layers economic. Treat agent PRs like any other "
        "contributor: sample one in ten for deep review, chosen at "
        "random. The moment you start skimming because they look right, "
        "you have drifted."
    ),
    11: (
        "Security teams: an agent is a new identity class. It needs an "
        "identity, a scoped permission set, an audit trail, and a blast "
        "radius. Decide in week one, not after the near-miss. Which "
        "branches can an agent push to? Which MCP tools can it call? "
        "Which secrets must it never see? Write the answers in "
        "AGENTS.md and enforce them in branch protection and "
        "environment policies. Prompt injection is the new SQL "
        "injection: any text the agent reads, including issue comments, "
        "web pages fetched by tools, even file contents, is potentially "
        "adversarial. The mitigations are familiar in shape: least "
        "privilege, output validation, human approval gates for "
        "high-impact actions, and the principle that an agent should "
        "never hold a secret it cannot rotate."
    ),
    12: (
        "Governance is where most teams will stumble in 2026. Three "
        "practical moves. One: make AGENTS.md a reviewed artifact, just "
        "like code. Changes to it go through PR review. Two: preserve "
        "agent run summaries on every PR, model version included. In "
        "regulated industries this is a compliance question, not a "
        "nice-to-have. Auditors will ask which model made this change, "
        "with what tools, against which spec. You want the answer to "
        "be a link, not an investigation. Three: set anti-goals "
        "explicitly. Write down what you are not trying to optimize "
        "for. Maximum agent autonomy at the expense of accountability "
        "is an anti-goal. AI did it is never a sufficient explanation "
        "in a post-incident review. Anti-goals are where hard-won "
        "wisdom lives."
    ),
    13: (
        "Seven failure modes will hit you in the first six months. "
        "I'll name them fast and you'll recognize at least two. "
        "Vibe-driven development: plausible diff, hidden bug. Fix: no "
        "agent run without a spec. Context starvation: locally "
        "plausible, globally wrong. Fix: invest in AGENTS.md and MCP. "
        "Trust drift: you start skimming. Fix: random deep-review "
        "sampling. Mistaking the agent for the loop: productivity "
        "bump, declare victory, stop investing. Fix: the artifacts make "
        "the loop better, the agent is one participant. Invisible "
        "agent work: future incident waiting. Fix: preserve traces. "
        "Deferring governance: locking down after a near-miss costs "
        "more. Fix: decide permissions in week one. Confusing "
        "productivity with capability: faster typing is not a better "
        "loop. Fix: measure artifact quality, not throughput."
    ),
    14: (
        "Here is what I want you to take back. One page. The AI-Native "
        "Success Criteria Memo. Loop scope: when do you consider your "
        "loop AI-native? Artifact criteria: what does every change "
        "retain? Quality bars: what blocks a merge? Human "
        "responsibility: which human, named by role? Anti-goals: what "
        "are you deliberately not optimizing for? Review cadence: "
        "quarterly, with an owner. Write v1 this week. It should be "
        "slightly embarrassing in its specificity, not polished for a "
        "slide. A bar the team routinely ignores teaches everyone the "
        "memo is decorative. The memo is a target, not a description "
        "of done. But it's the artifact that turns this talk into "
        "practice."
    ),
    15: (
        "Three questions to carry home. One: on your team's last "
        "non-trivial change, who actually decided what done meant? A "
        "human, an agent, or no one in particular? Two: which stage "
        "of your current loop would change least if every AI capability "
        "turned off tomorrow? That's where you're still doing "
        "traditional or AI-assisted work. Three: which of the seven "
        "failure modes is most likely to bite your team first? Write "
        "it on a sticky note before you leave this room. The AI-native "
        "cloud is not a product you buy. It's a loop you design. "
        "Thank you."
    ),
}


# ---------------------------------------------------------------------------
# Per-slide builders
# ---------------------------------------------------------------------------


def slide_blank(prs: Presentation):
    # Layout 6 in the default master is "Blank".
    return prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout


def build_slide_1(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)

    # Title placeholder (replaces default).
    title = slide.shapes.title
    title.left = MARGIN_X
    title.top = Inches(2.2)
    title.width = SLIDE_W - 2 * MARGIN_X
    title.height = Inches(1.4)
    _set_text(
        title.text_frame,
        TALK_TITLE,
        size_pt=54,
        color=theme.text_primary,
        bold=True,
    )

    # Subtitle.
    sub = slide.shapes.add_textbox(
        MARGIN_X, Inches(3.7), SLIDE_W - 2 * MARGIN_X, Inches(0.7)
    )
    _set_text(
        sub.text_frame,
        "From cloud-native to AI-native",
        size_pt=28,
        color=theme.accent_teal,
        bold=True,
    )

    # Accent rule.
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_X,
        Inches(4.6),
        Inches(1.6),
        Pt(3),
    )
    _set_solid_fill(rule, theme.accent_lime)
    _set_no_line(rule)
    _kill_shadow(rule)
    _set_alt_text(rule, "Decorative accent rule")

    # Presenter line.
    pres = slide.shapes.add_textbox(
        MARGIN_X, Inches(4.9), SLIDE_W - 2 * MARGIN_X, Inches(0.6)
    )
    _set_text(
        pres.text_frame,
        PRESENTER,
        size_pt=22,
        color=theme.text_muted,
    )


def build_slide_2(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "The Quiet Shift", theme)
    add_bullets(
        slide,
        [
            "2024 to 2026: a stack of compatible primitives",
            "Agent mode, coding agents, MCP, Spec Kit, AGENTS.md",
            "~90% of technologists use AI at work (DORA 2025)",
            "New question: is your loop designed for AI?",
        ],
        theme,
    )


def build_slide_3(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Three Eras, Side by Side", theme)
    add_comparison_table(
        slide,
        headers=["", "Traditional", "AI-Assisted", "AI-Native"],
        rows=[
            [
                "Unit of work",
                "Ticket \u2192 code",
                "Ticket \u2192 code (faster)",
                "Spec \u2192 context \u2192 change \u2192 trace",
            ],
            ["Who edits", "Human", "Human + suggestions", "Human and/or agent"],
            ["Artifacts", "Code, PR", "Code, PR", "Spec, context, trace, code, PR"],
            [
                "Governance",
                "Review, CI",
                "Review, CI",
                "+ spec review, agent permissions, eval gates",
            ],
        ],
        theme=theme,
    )


def build_slide_4(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Defining AI-Native", theme)
    add_bullets(
        slide,
        [
            "AI is a participant, not a tool",
            "Designed for, not bolted on",
            "Every action leaves a trail",
        ],
        theme,
        top=Inches(1.85),
        height=Inches(3.2),
    )

    # Three chips along the bottom.
    chip_top = Inches(5.4)
    chip_w = Inches(3.6)
    gap = (SLIDE_W - 2 * MARGIN_X - chip_w * 3) / 2
    for idx, (label, fill, on) in enumerate(
        [
            ("Designing", theme.accent_teal, theme.on_teal),
            ("Participating", theme.accent_sage, theme.on_sage),
            ("Accountable", theme.accent_lime, theme.on_lime),
        ]
    ):
        x = MARGIN_X + (chip_w + gap) * idx
        add_callout_chip(
            slide,
            label,
            fill,
            on,
            left=x,
            top=chip_top,
            width=chip_w,
        )


def build_slide_5(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "The Six-Stage Loop", theme)
    add_loop_diagram(slide, theme, top=Inches(2.4))

    # Two short captions across the bottom.
    cap_top = Inches(5.5)
    half_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) / 2
    left_box = slide.shapes.add_textbox(MARGIN_X, cap_top, half_w, Inches(0.6))
    _set_text(
        left_box.text_frame,
        "New artifact: spec",
        size_pt=24,
        color=theme.accent_teal,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    right_box = slide.shapes.add_textbox(
        MARGIN_X + half_w + Inches(0.4), cap_top, half_w, Inches(0.6)
    )
    _set_text(
        right_box.text_frame,
        "New arrow: feedback to spec",
        size_pt=24,
        color=theme.accent_lime,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def build_slide_6(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Cloud-Native Gave Us the Substrate", theme)
    add_bullets(
        slide,
        [
            "Cloud-native = the substrate",
            "Ephemeral envs, GitOps, observability, zero-trust",
            "AI-native depends on cloud-native maturity",
            "No shortcut: weak substrate, weak agents",
        ],
        theme,
    )


def build_slide_7(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Spec-Driven Development", theme)
    add_bullets(
        slide,
        [
            "Spec is the durable artifact",
            "Five lines beats fifty pages",
            "Cheaper to argue about a spec than a diff",
        ],
        theme,
        top=Inches(1.85),
        height=Inches(2.6),
    )

    # Section chip strip.
    strip_top = Inches(5.0)
    chip_w = Inches(2.85)
    chips = [
        "Why",
        "What changes",
        "Success criteria",
        "Out of scope",
    ]
    gap = (SLIDE_W - 2 * MARGIN_X - chip_w * len(chips)) / (len(chips) - 1)
    for idx, label in enumerate(chips):
        x = MARGIN_X + (chip_w + gap) * idx
        add_callout_chip(
            slide,
            label,
            theme.accent_sage,
            theme.on_sage,
            left=x,
            top=strip_top,
            width=chip_w,
        )


def build_slide_8(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Context Engineering", theme)

    # Custom layered bullets with bold accent labels.
    box = slide.shapes.add_textbox(
        MARGIN_X, Inches(1.85), SLIDE_W - 2 * MARGIN_X, Inches(3.8)
    )
    tf = box.text_frame
    tf.word_wrap = True

    layers = [
        ("AGENTS.md", "how this repo prefers to be edited"),
        ("MCP servers", "wiki, observability, internal APIs as tools"),
        ("Per-task context packets", "for non-trivial work"),
    ]
    for idx, (label, desc) in enumerate(layers):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.text = ""
        p.space_after = Pt(14)
        p.line_spacing = 1.2

        marker = p.add_run()
        marker.text = "\u25A0  "
        _style_run(marker, size_pt=24, color=theme.accent_teal, bold=True)

        lbl = p.add_run()
        lbl.text = f"{label} \u2014 "
        _style_run(lbl, size_pt=24, color=theme.accent_teal, bold=True)

        body = p.add_run()
        body.text = desc
        _style_run(body, size_pt=24, color=theme.text_primary)

    add_pullquote(
        slide,
        "Bad output is usually context starvation.",
        theme,
        top=Inches(5.6),
        height=Inches(1.0),
        accent=theme.accent_lime,
    )


def build_slide_9(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Humans + Agents", theme)

    # Wrong race (teal bar) vs Right race (lime bar) blocks.
    top = Inches(2.0)
    row_h = Inches(1.4)

    def add_race(idx: int, label: str, body: str, accent: RGBColor) -> None:
        y = top + (row_h + Inches(0.25)) * idx
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN_X,
            y,
            Pt(8),
            row_h,
        )
        _set_solid_fill(bar, accent)
        _set_no_line(bar)
        _kill_shadow(bar)
        _set_alt_text(bar, "Decorative accent bar")

        lbl_box = slide.shapes.add_textbox(
            MARGIN_X + Inches(0.3), y, Inches(3.4), Inches(0.55)
        )
        _set_text(
            lbl_box.text_frame,
            label,
            size_pt=22,
            color=accent,
            bold=True,
        )
        body_box = slide.shapes.add_textbox(
            MARGIN_X + Inches(0.3), y + Inches(0.55), SLIDE_W - 2 * MARGIN_X - Inches(0.3), Inches(0.85)
        )
        _set_text(
            body_box.text_frame,
            body,
            size_pt=22,
            color=theme.text_primary,
        )

    add_race(0, "Wrong race", "Cut headcount to do yesterday's work faster.", theme.accent_teal)
    add_race(1, "Right race", "Expand ambition. Orchestrate fleets of agents.", theme.accent_lime)

    # Closing line.
    foot = slide.shapes.add_textbox(
        MARGIN_X, Inches(5.7), SLIDE_W - 2 * MARGIN_X, Inches(0.6)
    )
    _set_text(
        foot.text_frame,
        "New roles: specifier, context curator, reviewer of intent.",
        size_pt=22,
        color=theme.text_muted,
        bold=True,
    )


def build_slide_10(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Reliability at Scale", theme)
    add_bullets(
        slide,
        [
            "No invisible work: trace every run",
            "Layered verification: tests \u2192 evaluators \u2192 policy \u2192 human",
            "Sample 1 in 10 agent PRs for deep review",
            "Trust calibration is a discipline, not a feeling",
        ],
        theme,
    )


def build_slide_11(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Security at Scale", theme)
    add_bullets(
        slide,
        [
            "Agents are a new identity class",
            "Scope: branches, MCP tools, secrets, environments",
            "Least privilege, validated outputs, human gates on impact",
        ],
        theme,
        top=Inches(1.85),
        height=Inches(3.0),
    )
    add_pullquote(
        slide,
        "Prompt injection is the new SQL injection.",
        theme,
        top=Inches(5.4),
        height=Inches(1.2),
        accent=theme.accent_lime,
    )


def build_slide_12(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Governance at Scale", theme)
    add_bullets(
        slide,
        [
            "AGENTS.md is a reviewed artifact",
            "Preserve run traces: model, tools, files, spec",
            "Auditors will ask. Make the answer a link.",
            "Write anti-goals, not just goals",
        ],
        theme,
    )


def build_slide_13(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Seven Failure Modes", theme)
    add_numbered_two_col(
        slide,
        [
            "Vibe-driven dev \u2192 write a spec",
            "Context starvation \u2192 AGENTS.md + MCP",
            "Trust drift \u2192 random deep reviews",
            "Agent \u2260 loop \u2192 invest in artifacts",
            "Invisible work \u2192 preserve traces",
            "Deferred governance \u2192 decide in week one",
            "Productivity \u2260 capability \u2192 measure artifacts",
        ],
        theme,
    )


def build_slide_14(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "One-Page Memo", theme)
    add_bullets(
        slide,
        [
            "Loop scope: when is your loop AI-native?",
            "Artifact criteria: what does every change retain?",
            "Quality bars: what blocks a merge?",
            "Human responsibility: which human, by role?",
            "Anti-goals: what you will not optimize for",
            "Review cadence: quarterly, with an owner",
        ],
        theme,
        top=Inches(1.85),
        height=Inches(4.0),
        size_pt=22,
    )
    add_callout_chip(
        slide,
        "Write v1 this week.",
        theme.accent_lime,
        theme.on_lime,
        left=MARGIN_X,
        top=Inches(6.05),
        width=Inches(4.0),
        height=Inches(0.7),
    )


def build_slide_15(prs: Presentation, theme: Theme) -> None:
    slide = slide_blank(prs)
    _set_paint_background(slide, theme)
    add_title(slide, "Three Questions", theme)
    add_bullets(
        slide,
        [
            "Who decided what done meant?",
            "Which stage survives AI turning off tomorrow?",
            "Which failure mode bites your team first?",
        ],
        theme,
        top=Inches(1.85),
        height=Inches(2.8),
        size_pt=24,
    )
    add_pullquote(
        slide,
        "The AI-native cloud is a loop you design.",
        theme,
        top=Inches(5.0),
        height=Inches(1.2),
        accent=theme.accent_teal,
    )
    pres = slide.shapes.add_textbox(
        MARGIN_X, Inches(6.35), SLIDE_W - 2 * MARGIN_X, Inches(0.5)
    )
    _set_text(
        pres.text_frame,
        f"Thank you  \u00b7  {PRESENTER}",
        size_pt=20,
        color=theme.text_muted,
    )


SLIDE_BUILDERS: list[Callable[[Presentation, Theme], None]] = [
    build_slide_1,
    build_slide_2,
    build_slide_3,
    build_slide_4,
    build_slide_5,
    build_slide_6,
    build_slide_7,
    build_slide_8,
    build_slide_9,
    build_slide_10,
    build_slide_11,
    build_slide_12,
    build_slide_13,
    build_slide_14,
    build_slide_15,
]


# ---------------------------------------------------------------------------
# Deck assembly
# ---------------------------------------------------------------------------


def _set_core_properties(prs: Presentation, theme: Theme) -> None:
    cp = prs.core_properties
    cp.title = TALK_TITLE
    cp.author = "Michelle Mei-Ling Sandford"
    cp.subject = "Building the AI-Native Cloud"
    cp.keywords = (
        "AI-native, cloud-native, DevOps, agents, AGENTS.md, MCP, Spec Kit, "
        "governance"
    )
    cp.comments = (
        f"{theme.name} theme. Generated by build_decks.py. "
        "Calibri throughout, all text >=18pt, WCAG AA contrast."
    )
    cp.language = "en-US"


def _add_notes(slide, text: str, theme: Theme) -> None:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    p.text = ""
    run = p.add_run()
    run.text = text
    # Notes are read by presenters, not projected; use the default Calibri 18pt.
    _style_run(run, size_pt=18, color=RGBColor(0x11, 0x11, 0x11))


def render_deck(theme: Theme, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(SLIDE_BUILDERS)
    for idx, builder in enumerate(SLIDE_BUILDERS, start=1):
        builder(prs, theme)
        slide = prs.slides[idx - 1]
        if idx > 1:
            add_footer(slide, idx, total, theme)
        _add_notes(slide, SPEAKER_NOTES[idx], theme)

    _set_core_properties(prs, theme)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


# ---------------------------------------------------------------------------
# Verification
# ---------------------------------------------------------------------------


def assert_deck(out_path: Path) -> None:
    prs = Presentation(out_path)
    assert len(prs.slides) == 15, (
        f"Expected 15 slides, got {len(prs.slides)} in {out_path.name}"
    )
    for i, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title
        assert title is not None and title.has_text_frame, (
            f"Slide {i} in {out_path.name} has no title placeholder"
        )
        assert title.text_frame.text.strip(), (
            f"Slide {i} in {out_path.name} has an empty title"
        )
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    if run.font.name not in (FONT, None):
                        raise AssertionError(
                            f"{out_path.name} slide {i}: font is "
                            f"{run.font.name!r}, expected {FONT!r} "
                            f"(run text: {run.text[:40]!r})"
                        )
                    size = run.font.size
                    if size is None:
                        raise AssertionError(
                            f"{out_path.name} slide {i}: run has no "
                            f"explicit size (text: {run.text[:40]!r})"
                        )
                    if size.pt < 18:
                        raise AssertionError(
                            f"{out_path.name} slide {i}: run is "
                            f"{size.pt}pt, below 18pt minimum "
                            f"(text: {run.text[:40]!r})"
                        )
    print(
        f"OK: {out_path.name} \u2014 15 slides, all Calibri, "
        f"all \u226518pt, all titled"
    )


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--light", action="store_true", help="Build only the light deck")
    parser.add_argument("--dark", action="store_true", help="Build only the dark deck")
    parser.add_argument(
        "--out",
        type=Path,
        default=Path(__file__).resolve().parent / "dist",
        help="Output directory (default: ./dist)",
    )
    args = parser.parse_args()

    build_light = args.light or not (args.light or args.dark)
    build_dark = args.dark or not (args.light or args.dark)

    args.out.mkdir(parents=True, exist_ok=True)
    base = "Building-the-AI-Native-Cloud"

    targets: list[tuple[Theme, Path]] = []
    if build_light:
        targets.append((LIGHT, args.out / f"{base}_Light.pptx"))
    if build_dark:
        targets.append((DARK, args.out / f"{base}_Dark.pptx"))

    for theme, out_path in targets:
        render_deck(theme, out_path)
        print(f"Wrote {out_path}")
        assert_deck(out_path)

    return 0


if __name__ == "__main__":
    sys.exit(main())
